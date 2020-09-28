import pandas as pd
import yaml
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement


# Glasswall palette
dark_blue = RGBColor(14, 61, 90)
green_blue = RGBColor(26, 145, 154) 
white = RGBColor(255, 255, 255) 
blue1 = RGBColor(22, 94, 122) 
# table colors:
gray = RGBColor(191, 191, 191)
blue2 = RGBColor(45, 92, 117)


# Letter font
gw_font = 'Mediator Narrow'
#gw_font = 'Asterisk Sans Pro Regular'
#gw_font = 'Liberation Sans'
#gw_font = 'Calibry' 
#gw_font = 'Times New Roman'
#gw_font = 'Atwic Regular Reversed'
#gw_font = 'Umba Soft Alt Regular'


def examine_template():

    prs = Presentation()
    
    for n in range(0, 11):

        slide = prs.slides.add_slide(prs.slide_layouts[n])
        
        print('Master Slide ' + str(n))

        for shape in slide.placeholders:
            print('%d, %s' % (shape.placeholder_format.idx, shape.name))
 
 
def logo(slide, img_path='images/glasswall_logo.png', place='top'):
   
    if place == 'top':
        # Logo size
        width = Inches(1.2)
        height = Inches(0.6) # width half

        # Place it on top right corner
        top = Inches(0.1) 
        left = Inches(10) - width - Inches(0.2)
        
    elif place == 'center':
        # Logo size
        width = Inches(6)
        height = Inches(3) # width half

        # Place it in the center
        left = (Inches(10) - width)/2
        top = (Inches(7.5) - height)/2
    
    pic = slide.shapes.add_picture(img_path, left, top, width, height) 
    
    
def set_background_color(slide, bg_color=dark_blue):
    background = slide.background  
    fill(background, bg_color)
 

def fill(shape, fill_color=dark_blue):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
               
    
def make_presentation(sheet_name, output_to):

    prs = Presentation()

    # PROJECT SLIDES
    
    if sheet_name == 'Projects Team Structure':
    
        with open(os.path.join('yaml_files', sheet_name+'.yml'), 'r') as f:
            df = pd.json_normalize(yaml.safe_load(f, Loader=yaml.FullLoader))
        
        with open('yaml_files/Resource and responsability.yml', 'r') as f:
            df2 = pd.json_normalize(yaml.safe_load(f, Loader=yaml.FullLoader))
            
        for row_index in range(len(df)):
            add_project_slide(prs, df, row_index, df2)

    prs.save(output_to)
    
    
def add_project_slide(prs, df, row_index, df2):

    p_name = df.iloc[row_index]['project_name'].rstrip('\n')
    p_priority = df.iloc[row_index]['priority']
    
    p_stakeholder = df.iloc[row_index]['business_owner_stakeholders'].split('\n')
    p_stakeholder = list(filter(('').__ne__, p_stakeholder))
    
    p_manager = df.iloc[row_index]['delivery_manager'].replace(',', '\n').split('\n')
    p_manager = list(filter(('').__ne__, p_manager))
    
    p_security = df.iloc[row_index]['security_champ'].split('\n')
    p_security = list(filter(('').__ne__, p_security))

    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    
    set_background_color(slide)
    
    logo(slide)
    
    shapes = slide.shapes
    
    # TITLE
    title = shapes.title
    
    p = ''
    if '1' in p_priority:
        p = '1'
    elif '2' in p_priority:
        p = '2'
    elif '3' in p_priority:
        p = '3'
    
    title.text = 'PROJECT: ' + p_name.upper() + '\n Priority: ' + p
    
    text_settings(title, i=0, font_size=Pt(24))    
    text_settings(title, i=1, font_size=Pt(22), font_color=green_blue)


    # ROW OF RECTANGLES

    # RECTANGLE 1        
    add_rectangle(shapes, p_stakeholder, 'Stakeholder', blue1, white, left=Inches(0.75))
    
    # RECTANGLE 2
    add_rectangle(shapes, p_manager, 'Delivery Manager', green_blue, white, left=Inches(3.75))

    # RECTANGLE 3         
    add_rectangle(shapes, p_security, 'Security Champ', white, dark_blue, left=Inches(6.75))
        
   # TABLE 
    rnr = df2[df2['project_name']==p_name].reset_index()
    
    if len(rnr)>0:
        if rnr['project_name'][0] == '#Security-Privacy-Champion': 
            add_table(shapes, rnr, blue1, top=Inches(2.0))
        elif rnr['project_name'][0] == 'GlassWall Resources':
            rnr1 = rnr.iloc[:9]
            rnr2 = rnr.iloc[9:,:].reset_index()
            
            add_table(shapes, rnr1, blue1, Inches(2.0), col_width=Inches(2.25), left=Inches(0.5))
            add_table(shapes, rnr2, blue1, Inches(2.0), col_width=Inches(2.25), left=Inches(5.25))
        else:
            add_table(shapes, rnr, blue1)
            
            
def text_settings(shape, i=0, aligment=PP_ALIGN.LEFT, font_color=white, font_size=Pt(12), font=gw_font, bold=False):
    text = shape.text_frame.paragraphs[i]
    text.alignment = aligment
    text.font.name = font
    text.font.size = font_size
    text.font.color.rgb = font_color
    text.font.bold = bold
    
    
def add_rectangle(shapes, names, rol, fill_color, font_color, left, width=Inches(2.5)):

    n = len(names)
    if n>0:
        sentence = names[0]
        for i in range(1,n):
            sentence = sentence + '\n' + names[i]      
        sentence = sentence
        
        top = Inches(2.25)    
        height = Inches(0.65)    
        if n>1:
            top -= Inches(n*0.15)
            height += Inches(n*0.15)
            
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.text = sentence

        for i in range(n):
            text_settings(shape, i, PP_ALIGN.CENTER, font_size=Pt(14))
                  
        fill(shape)
        
        line = shape.line
        line.color.rgb = white
        line.width = Pt(1.5)
        
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left+Inches(0.25), top-Inches(0.25), width-Inches(0.5), height=Inches(0.4))
        shape.text = rol.upper()
                
        text_settings(shape, aligment=PP_ALIGN.CENTER, font_color=font_color)
        
        fill(shape, fill_color)
        
        line = shape.line
        line.color.rgb = fill_color
    
    
def add_table(shapes, df, table_color, top=Inches(3.5), col_width=Inches(4.0), left=Inches(1.0), width=Inches(6.0), height=Inches(0.8)):
             
     cols = 2
     rows = len(df)+1
    
     shape = shapes.add_table(rows, cols, left, top, width, height)
     table = shape.table
     
     # set column widths
     table.columns[0].width = col_width
     table.columns[1].width = col_width

     # write column headings
     table.cell(0, 0).text = 'Resource'.upper()    
     table.cell(0, 1).text = 'Responsability'.upper()

     # write body cells
     for i in range(1, rows):
        table.cell(i, 0).text = df['resource_name'][i-1]
        
        cell = table.cell(i, 0)
        fill(cell, blue2)
        text_settings(cell, aligment=PP_ALIGN.CENTER)
        set_cell_border(cell, blue2, white)
                     
        table.cell(i, 1).text = df['responsibility_for_the_project'][i-1] 
        
        cell = table.cell(i, 1)
        fill(cell, blue2)    
        text_settings(cell, aligment=PP_ALIGN.CENTER)
        set_cell_border(cell, blue2, white)
        
     # set headings color and font
     cell = table.cell(0, 0)
     fill(cell, gray)   
     text_settings(cell, aligment=PP_ALIGN.CENTER, font_color=blue2, bold=True)    
     set_cell_border(cell, gray, gray)
    
     cell = table.cell(0, 1)
     fill(cell, gray)    
     text_settings(cell, aligment=PP_ALIGN.CENTER, font_color=blue2, bold=True)    
     set_cell_border(cell, gray, gray)
        

def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element
        

def set_cell_border(cell, border_color_LR, border_color_TB, border_width='12700'):

    # convert RGB to hex
    border_color_LR = '%02x%02x%02x' % border_color_LR
    border_color_TB = '%02x%02x%02x' % border_color_TB
    
    colors = [border_color_LR, border_color_LR, border_color_TB, border_color_TB]
    
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    
    lines = ['a:lnL','a:lnR','a:lnT','a:lnB']
    
    for line, color in zip(lines,  colors):
        ln = SubElement(tcPr, line, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med') 
        