import pandas as pd
import yaml
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
#import subprocess as sp


# Include this path if working in Google Colab    
#d = '/content/gdrive/My Drive/p2-p-data-visualization-Pappaterra-Lucia/'
# If not 
d = ''


# Glasswall palette
dark_blue = RGBColor(14, 61, 90)
green_blue = RGBColor(26, 145, 154) 
white = RGBColor(255, 255, 255) 
blue1 = RGBColor(22, 94, 122) 
# table colors:
gray = RGBColor(191, 191, 191)
blue2 = RGBColor(45, 92, 117)
#rag colors:
green = RGBColor(0, 204, 153) #(0, 255, 0)
amber = RGBColor(255, 204, 0) #(255, 153, 51)
red = RGBColor(255, 102, 102) #(255, 0, 0)

# Letter font
gw_font = 'Lato'


def examine_template():

    prs = Presentation()
    
    for n in range(0, 11):

        slide = prs.slides.add_slide(prs.slide_layouts[n])
        
        print('Master Slide ' + str(n))

        for shape in slide.placeholders:
            print('%d, %s' % (shape.placeholder_format.idx, shape.name))
 
 
def logo(slide, img_path=d+'images/glasswall_logo.png', place='top right'):
   
    if place == 'top right':
        # Logo size
        width = Inches(1.2)
        height = Inches(0.6) # width half

        # Place it on top right corner
        top = Inches(0.1) 
        left = Inches(10.0) - width - Inches(0.2)
        
    elif place == 'center':
        # Logo size
        width = Inches(6.0)
        height = Inches(3.0) # width halfs

        # Place it in the center
        left = (Inches(10.0) - width)/2
        top = (Inches(7.5) - height)/2
        
    elif place == 'top left':
        # Logo size
        width = Inches(1.25)
        height = Inches(0.625) # width half

        # Place it on top left corner
        top = Inches(0.25) 
        left = Inches(0.3) 
    
    pic = slide.shapes.add_picture(img_path, left, top, width, height) 
    
    
def set_background_color(slide, bg_color=dark_blue):
    background = slide.background  
    fill(background, bg_color)
 

def fill(shape, fill_color=dark_blue):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
    
    
def sheet_to_dfs(sheet_name):

    if sheet_name == 'Projects Team Structure':
    
        with open(os.path.join(d+'yaml_files', sheet_name+'.yml'), 'r') as f:
            df1 = pd.json_normalize(yaml.safe_load(f))
            #df1 = pd.json_normalize(yaml.load(f, Loader=yaml.FullLoader))
                    
        with open(d+'yaml_files/Resource and responsability.yml', 'r') as f:
            df2 = pd.json_normalize(yaml.safe_load(f))
            #df2 = pd.json_normalize(yaml.load(f, Loader=yaml.FullLoader))
            
        return df1, df2
               
    
def make_presentation(sheet_name, output_to, single=False, dm=False):

    prs = Presentation()
     
    if sheet_name == 'Projects Team Structure':
    
        df1, df2 = sheet_to_dfs(sheet_name)
        
        # PROJECT SLIDES    
        for row_index in range(len(df1)):
            add_project_slide(prs, df1, row_index, df2)
            
            if single:           
                # Also save single presentations to single folder
                prs2 = Presentation()           
                add_project_slide(prs2, df1, row_index, df2)
                #p_name = df1.iloc[row_index]['project_name'].rstrip('\n') 
                p_slack_channel = df1.iloc[row_index]['slack_channel'][1:]
                prs2.save(d+'outputs/single/' + p_slack_channel + '.pptx')
                
                # also save it as pdf
                #sp.call(['libreoffice', '--headless', '--convert-to', 'pdf', p_name], cwd='outputs/single')
                
        if dm:
            # Also create a presentation per delivery manager to dm folder       
            d_managers = list(df1['delivery_manager'].unique())
            
            for elem in d_managers:
                if ',' in elem:
                    d_managers.remove(elem)
                       
            for d_manager in d_managers:
                prs3 = Presentation()
                
                df3 = df1[df1['delivery_manager'].str.contains(d_manager)]
                df3 = df3.reset_index(drop=True)
                
                # PROJECT SLIDES    
                for row_index in range(len(df3)):
                    add_project_slide(prs3, df3, row_index, df2)
                       
                prs3.save(d+'outputs/DM/' + d_manager + '.pptx')
                
                # also save it as pdf
                #sp.call(['libreoffice', '--headless', '--convert-to', 'pdf', d_manager], cwd='outputs/DM')

    prs.save(output_to)
    
    # also save it as pdf
    #sp.call(['libreoffice', '--headless', '--convert-to', 'pdf', output_to.split('/')[1]], cwd='outputs')
 

def add_project_slide(prs, df, row_index, df2):

    p_name = df.iloc[row_index]['project_name']
    p_priority = df.iloc[row_index]['priority']
    p_slack_channel = df.iloc[row_index]['slack_channel']
    p_description = df.iloc[row_index]['short_project_description']
    
    p_stakeholder = df.iloc[row_index]['stakeholders'].split('\n')
    #filter empty string from list:
    p_stakeholder = list(filter(('').__ne__, p_stakeholder)) 
    
    p_manager = df.iloc[row_index]['delivery_manager'].split(', ')
    p_manager = list(filter(('').__ne__, p_manager))
    
    p_security = df.iloc[row_index]['security_champ'].split('\n')
    p_security = list(filter(('').__ne__, p_security))  

    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    
    set_background_color(slide)
    
    logo(slide)
    
    shapes = slide.shapes
    
    # TITLE
    title = shapes.title
    
    p = ''
    if '1' in p_priority:
        p = '1'
    elif '2' in p_priority:
        p = '2'
    elif '3' in p_priority:
        p = '3'
    elif 'ByRequest' in p_priority:
        p = 'by request'
    
    title.text = '\n' + '\n PROJECT: ' + p_name.upper()  + '\n' + p_description + '\n Priority: ' + p
    text_settings(title, i=0)
    text_settings(title, i=1)
    text_settings(title, i=2, font_size=Pt(26))
    
    description_font = Pt(12)
    if len(p_description)>400: description_font = Pt(10)    
    text_settings(title, i=3, font_size=description_font)
    
    text_settings(title, i=4, font_size=Pt(24), font_color=green_blue)

    # ROW OF RECTANGLES

    # RECTANGLE 1        
    add_rectangle(shapes, p_stakeholder, 'Stakeholder', blue1, white, left=Inches(0.75))
    
    # RECTANGLE 2
    add_rectangle(shapes, p_manager, 'Delivery Manager', green_blue, white, left=Inches(3.75))

    # RECTANGLE 3         
    add_rectangle(shapes, p_security, 'Security Champ', white, dark_blue, left=Inches(6.75))
        
   # TABLE 
    rnr = df2[(df2['slack_channel'] == p_slack_channel)].reset_index()
    
    if len(rnr)>0:
        add_table(shapes, rnr, blue1)
            
            
def text_settings(shape, i=0, alignment=PP_ALIGN.LEFT, font_color=white, font_size=Pt(14), font=gw_font, bold=False):
    text = shape.text_frame.paragraphs[i]
    text.alignment = alignment
    text.font.name = font
    text.font.size = font_size
    text.font.color.rgb = font_color
    text.font.bold = bold
    
    
def add_rectangle(shapes, names, rol, fill_color, font_color, left, width=Inches(2.5)):

    n = len(names)
    if n>0:
        sentence = names[0]
        for i in range(1,n):
            sentence = sentence + '\n' + names[i]      
        sentence = sentence
        
        top = Inches(2.65)    
        height = Inches(0.65)    
        if n>1:
            top -= Inches(n*0.15)
            height += Inches(n*0.15)
            
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.text = sentence

        for i in range(0, n):
            text_settings(shape, i, PP_ALIGN.CENTER, font_size=Pt(16))
                  
        fill(shape)
        
        line = shape.line
        line.color.rgb = white
        line.width = Pt(1.5)
        
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left+Inches(0.25), top-Inches(0.25), width-Inches(0.5), height=Inches(0.4))
        shape.text = rol.upper()
                
        text_settings(shape, alignment=PP_ALIGN.CENTER, font_color=font_color, font_size=Pt(12))
        
        fill(shape, fill_color)
        
        line = shape.line
        line.color.rgb = fill_color
    
    
def add_table(shapes, df, table_color, top=Inches(3.65), col_width=Inches(4.0), left=Inches(1.0), width=Inches(6.0), height=Inches(0.8)):
             
     cols = 2
     rows = len(df)+1
    
     shape = shapes.add_table(rows, cols, left, top, width, height)
     table = shape.table
     
     # set column widths
     table.columns[0].width = col_width
     table.columns[1].width = col_width

     # write column headings
     table.cell(0, 0).text = 'Resource'.upper()    
     table.cell(0, 1).text = 'Responsability'.upper()

     # write body cells
     for i in range(1, rows):
        table.cell(i, 0).text = df['resource_name'][i-1]
        
        cell = table.cell(i, 0)
        fill(cell, blue2)
        text_settings(cell, alignment=PP_ALIGN.CENTER)
        set_cell_border(cell, blue2, white)
                     
        table.cell(i, 1).text = df['resource_responsability'][i-1] 
        
        cell = table.cell(i, 1)
        fill(cell, blue2)    
        text_settings(cell, alignment=PP_ALIGN.CENTER)
        set_cell_border(cell, blue2, white)
        
     # set headings color and font
     cell = table.cell(0, 0)
     fill(cell, gray)   
     text_settings(cell, alignment=PP_ALIGN.CENTER, font_color=blue2, bold=True)    
     set_cell_border(cell, gray, gray)
    
     cell = table.cell(0, 1)
     fill(cell, gray)    
     text_settings(cell, alignment=PP_ALIGN.CENTER, font_color=blue2, bold=True)    
     set_cell_border(cell, gray, gray)
        

def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element
        

def set_cell_border(cell, border_color_LR, border_color_TB, border_width='12700'):

    # convert RGB to hex
    border_color_LR = '%02x%02x%02x' % border_color_LR
    border_color_TB = '%02x%02x%02x' % border_color_TB
    
    colors = [border_color_LR, border_color_LR, border_color_TB, border_color_TB]
    
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    
    lines = ['a:lnL','a:lnR','a:lnT','a:lnB']
    
    for line, color in zip(lines,  colors):
        ln = SubElement(tcPr, line, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med') 
        
        
def daily_presentation(df, output_to):

    prs = Presentation()
    
    for row_index in range(len(df)):
        add_slide(prs, df, row_index)

    prs.save(output_to)
    
    # also save it as pdf
    #import subprocess as sp
    #sp.call(['libreoffice', '--headless', '--convert-to', 'pdf', output_to.split('/')[1]], cwd='outputs')
    
    
def add_slide(prs, df, row_index):

    # Project information
    priority = df.iloc[row_index]['Priority']
    service_name = df.iloc[row_index]['Service Name']
    date = df.iloc[row_index]['Date']
    p_name = df.iloc[row_index]['Project']
    d_manager = df.iloc[row_index]['Delivery Manager']
    onboarding = df.iloc[row_index]['Onboarding']
    offboarding = df.iloc[row_index]['Offboarding']
    in_progress = df.iloc[row_index]['Issues in progress']
    closed = df.iloc[row_index]['Closed issues']
    major_issues = df.iloc[row_index]['Major issues'] 
    rag_status = df.iloc[row_index]['RAG Status']
    team = df.iloc[row_index]['Team members']
    next_release = df.iloc[row_index]['Next Release/ Important dates']
    comments = df.iloc[row_index]['If you have any more comments']
    parent_service = df.iloc[row_index]['Parent Service']
    service_manager = df.iloc[row_index]['Service Manager']
    service_name = df.iloc[row_index]['Service Name']
    
    if len(next_release)>250 and comments == '':
        comments = next_release[len(next_release)//2:]
        next_release = next_release[:len(next_release)//2]
    
    if rag_status == 'Green':
        rgb_color = green
    elif rag_status == 'Red':
        rgb_color = red
    elif rag_status == 'Amber':
        rgb_color = amber

    #slide
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
        
    logo(slide, img_path=d+'images/glasswall_logo2.png', place='top left')
    
    shapes = slide.shapes
    
    # TITLE
    title = shapes.title
    
    if service_name != '':
        project_name = service_name
    else:
        project_name = p_name   
    
    title.text = date +'\n' + project_name + '\n\n'
    text_settings(title, i=0, alignment=PP_ALIGN.CENTER, font_size=Pt(12), font_color=dark_blue, bold=True)  
    text_settings(title, i=1, alignment=PP_ALIGN.CENTER, font_size=Pt(24), font_color=dark_blue, bold=True)    
    
    # FIGURES MEASURES:
    
    # distance to left side
    left = Inches(0.35) #1.25, 1.1, 0.8, 0.65, 0.35
    
    # distance to top
    top = Inches(1.1)
    
    # widths:
    width0 = Inches(1.0) # Priority
    width1 = Inches(0.9)
    width2 = Inches(1.4) # RAG Status
    width3 = Inches(1.1)
    width4 = Inches(1.9) # Bussines Owner
    width5 = Inches(3.0)
    
    width8 = Inches(10.0) - 2*left # Issues
    width9 = width8/3 # In progress, Closed, Major Issues
    
    # Team, Onboarding, Offboarding, Next Release, Comments:
    width10 = Inches(2.0) 
    width11 = Inches(10.0) - 2*left - width10
    width12 = (width11 - width10)/2
    
    # heights:
    height0 = Inches(0.3) 
    height1 = Inches(2.95)
    
    # lefts:
    left1 = left + width10  
    left2 = left1 + width12
    left3 = left1 + width12 + width10
    
    # separation between figures:
    sep0 = Inches(0.1)
    sep1 = Inches(0.03)
      
    
    # First line
    make_rectangle(shapes, 'Priority:', left, top, width0, height0)
    left_aux = left + width0
    make_rectangle(shapes, priority, left_aux, top, width1, height0, font_color=dark_blue, fill_color=white, line_color=gray, alignment=PP_ALIGN.CENTER, bold=True)
    
    left_aux = left_aux + width1
    make_rectangle(shapes, 'RAG Status:', left_aux, top, width2, height0)   
    left_aux = left_aux + width2
    make_rectangle(shapes, rag_status, left_aux, top, width3, height0, font_color=dark_blue, fill_color=rgb_color, line_color=rgb_color, alignment=PP_ALIGN.CENTER, bold=True)
    
    left_aux = left_aux + width3
    make_rectangle(shapes, 'Service Manager', left_aux, top, width4, height0)
    left_aux = left_aux + width4 
    make_rectangle(shapes, service_manager, left_aux, top, width5, height0, font_color=dark_blue, fill_color=white, line_color=gray)
    
    # Next line
    top = top + height0 + sep0 
    
    make_rectangle(shapes, 'Issues', left, top, width8, height0, alignment=PP_ALIGN.CENTER, fill_color=blue2)
    
    # Next line
    top = top + height0 + sep1    
    
    make_rectangle(shapes, 'In Progress', left, top, width9, height0, alignment=PP_ALIGN.CENTER)
    make_rectangle(shapes, 'Closed', left + width9, top, width9, height0, alignment=PP_ALIGN.CENTER)
    make_rectangle(shapes, 'Major Issues', left + 2*width9, top, width9, height0, alignment=PP_ALIGN.CENTER)    
    
    # Next line
    top = top + height0 + sep1
    
    make_rectangle(shapes, in_progress, left, top, width9, height1, font_color=green_blue, bold=True, fill_color=white, line_color=gray)
    make_rectangle(shapes, closed, left + width9, top, width9, height1, font_color=green_blue, bold=True, fill_color=white, line_color=gray)
    make_rectangle(shapes, major_issues, left + 2*width9, top, width9, height1, font_color=green_blue, bold=True, fill_color=white, line_color=gray)
    
    # Next line
    top = top + height1 + sep0
      
    make_rectangle(shapes, 'Delivery Manager:', left, top, width10, height0)
    make_rectangle(shapes, d_manager, left1, top, width12, height0, font_color=dark_blue, fill_color=white, line_color=gray) 

    make_rectangle(shapes, 'Parent Service:', left2, top, width10, height0) 
    make_rectangle(shapes, parent_service, left3, top, width12, height0, font_color=dark_blue, fill_color=white, line_color=gray)    

    # Next line 
    top = top + height0 + sep0  

    make_rectangle(shapes, 'Slack Channel:', left, top, width10, height0) 
    
    if ' ' not in p_name:
        slack_channel = '#'+p_name
    else:
        slack_channel = ''
    
    make_rectangle(shapes, slack_channel, left1, top, width12, height0, font_color=dark_blue, fill_color=white, line_color=gray)
    
    make_rectangle(shapes, 'Next Release:', left2, top, width10, height0) 
    make_rectangle(shapes, next_release, left3, top, width12, height0, font_color=green_blue, bold=True, fill_color=white, line_color=gray)
  
    # Next line 
    top = top + height0 + sep0  
       
    make_rectangle(shapes, 'Team:', left, top, width10, height0) 
    make_rectangle(shapes, team, left1, top, width11, height0, font_color=dark_blue, fill_color=white, line_color=gray)
    
    # Next line
    top = top + height0 + sep0
    
    if onboarding == 'None':
        color, b = dark_blue, False
    else:
        color, b = green_blue, True
      
    make_rectangle(shapes, 'Onboarding:', left, top, width10, height0)  
    make_rectangle(shapes, onboarding, left1, top, width12, height0, font_color=color, bold=b, fill_color=white, line_color=gray) 
    
    if offboarding == 'None':
        color, b = dark_blue, False
    else:
        color, b = green_blue, True
         
    make_rectangle(shapes, 'Offboarding:', left2, top, width10, height0) 
    make_rectangle(shapes, offboarding, left3, top, width12, height0, font_color=color, bold=b, fill_color=white, line_color=gray)
        
    # Next line
    top = top + height0 + sep0
    
    make_rectangle(shapes, 'Comments:', left, top, width10, height0) 
    make_rectangle(shapes, comments, left1, top, width11, height0, font_color=green_blue, bold=True, fill_color=white, line_color=gray)

    
def make_rectangle(shapes, sentence, left, top, width, height, font_color=white, fill_color=dark_blue, line_color=dark_blue, alignment=PP_ALIGN.LEFT, bold=False, font_size=Pt(14)):
    
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = sentence
    
    if height >= Inches(2.0):
        if len(shape.text)>550:
            font_size=Pt(10)          
        elif len(shape.text)>=500:  
            font_size=Pt(11)
        elif len(shape.text)>=350:
            font_size=Pt(12)
        elif len(shape.text)>200: 
            font_size=Pt(13)
    else: 
        if len(shape.text)>100:
            font_size=Pt(11)   
    
    text_settings(shape, font_color=font_color, alignment=alignment, bold=bold, font_size=font_size)
    
    fill(shape, fill_color=fill_color)
    line = shape.line
    line.color.rgb = line_color
       